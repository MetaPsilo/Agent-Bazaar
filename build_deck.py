#!/usr/bin/env python3
"""Build AgentBazaar PowerPoint pitch deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import os
import shutil

# Constants
BG = RGBColor(9, 9, 11)
WHITE = RGBColor(250, 250, 250)
ACCENT = RGBColor(59, 130, 246)
SECONDARY = RGBColor(161, 161, 170)
MUTED = RGBColor(113, 113, 122)
FONT = 'Calibri'
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BASE = '/Users/dan/projects/agent-bazaar'
ASSETS = os.path.join(BASE, 'pitch-assets')
LOGO = os.path.join(BASE, 'frontend/public/logo.png')

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_text(slide, left, top, width, height, text, size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tf

def add_para(tf, text, size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = align
    return p

def add_bullet(tf, text, size=20, color=RGBColor(212, 212, 216)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = FONT
    p.space_before = Pt(6)
    p.level = 0
    # Add blue bullet via run
    p.clear()
    run = p.add_run()
    run.text = "●  "
    run.font.size = Pt(10)
    run.font.color.rgb = ACCENT
    run.font.name = FONT
    run2 = p.add_run()
    run2.text = text
    run2.font.size = Pt(size)
    run2.font.color.rgb = color
    run2.font.name = FONT
    return p

def add_numbered(tf, num, text, size=20, color=RGBColor(212, 212, 216)):
    p = tf.add_paragraph()
    p.clear()
    p.space_before = Pt(6)
    run = p.add_run()
    run.text = f"{num}.  "
    run.font.size = Pt(size)
    run.font.color.rgb = ACCENT
    run.font.bold = True
    run.font.name = FONT
    run2 = p.add_run()
    run2.text = text
    run2.font.size = Pt(size)
    run2.font.color.rgb = color
    run2.font.name = FONT

def divider(slide, left, top, width=Inches(0.8)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

def add_box(slide, left, top, w, h, text, subtitle="", border_color=ACCENT, bg_color=RGBColor(17,17,19)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = FONT
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(10)
        p.font.color.rgb = MUTED
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
    return shape

def arrow_right(slide, left, top, width=Inches(0.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

def slide_number(slide, num):
    add_text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4), str(num), size=11, color=MUTED, align=PP_ALIGN.RIGHT)

# ─── SLIDE 1: Title ───
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)
# Logo
logo_w = Inches(1.8)
logo_left = (SLIDE_W - logo_w) // 2
slide.shapes.add_picture(LOGO, logo_left, Inches(1.8), logo_w)
# Title
tf = add_text(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(2), "", size=40, bold=True, align=PP_ALIGN.CENTER)
p = tf.paragraphs[0]
p.clear()
r1 = p.add_run()
r1.text = "The Permissionless Protocol for "
r1.font.size = Pt(40)
r1.font.bold = True
r1.font.color.rgb = WHITE
r1.font.name = FONT
r2 = p.add_run()
r2.text = "AI Agent Commerce"
r2.font.size = Pt(40)
r2.font.bold = True
r2.font.color.rgb = ACCENT
r2.font.name = FONT
p.alignment = PP_ALIGN.CENTER
# URL
add_text(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5), "agentbazaar.org", size=18, color=ACCENT, align=PP_ALIGN.CENTER)
slide_number(slide, 1)

# ─── SLIDE 2: The Problem ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
divider(slide, Inches(1.2), Inches(1.2))
add_text(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(1), "The Problem", size=44, bold=True)
tf = add_text(slide, Inches(1.2), Inches(2.7), Inches(10), Inches(1), "", size=22)
p = tf.paragraphs[0]
p.clear()
r = p.add_run(); r.text = "AI agents are exploding — but they have "; r.font.size = Pt(22); r.font.color.rgb = SECONDARY; r.font.name = FONT
r = p.add_run(); r.text = "no way to buy or sell services"; r.font.size = Pt(22); r.font.color.rgb = ACCENT; r.font.bold = True; r.font.name = FONT
r = p.add_run(); r.text = " from each other."; r.font.size = Pt(22); r.font.color.rgb = SECONDARY; r.font.name = FONT

tf2 = add_text(slide, Inches(1.2), Inches(3.8), Inches(10), Inches(3.5), "", size=20)
add_bullet(tf2, "No discovery mechanism")
add_bullet(tf2, "No autonomous payments")
add_bullet(tf2, "No reputation system")
add_bullet(tf2, "Every agent platform is centralized and siloed")
slide_number(slide, 2)

# ─── SLIDE 3: The Opportunity ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
divider(slide, Inches(1.2), Inches(1.0))
add_text(slide, Inches(1.2), Inches(1.3), Inches(10), Inches(0.8), "The Opportunity", size=44, bold=True)
tf = add_text(slide, Inches(1.2), Inches(2.3), Inches(5.5), Inches(4), "", size=20)
add_bullet(tf, "AI agent market projected to reach $47B+ by 2030")
add_bullet(tf, "Agents need to transact autonomously — can't swipe a credit card")
add_bullet(tf, 'Winner takes the "commerce layer" for the agent economy')
add_bullet(tf, "Think: what the App Store did for mobile apps, but for AI agents")

# Bar chart
chart_data = CategoryChartData()
chart_data.categories = ['2024', '2025', '2026', '2027', '2030']
chart_data.add_series('Market Size ($B)', (5, 12, 22, 33, 47))
chart_left = Inches(7.2)
chart_top = Inches(2.0)
chart_w = Inches(5.5)
chart_h = Inches(4.5)
chart_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_left, chart_top, chart_w, chart_h, chart_data)
chart = chart_frame.chart
chart.has_legend = False
# Style the chart
plot = chart.plots[0]
series = plot.series[0]
series.format.fill.solid()
series.format.fill.fore_color.rgb = ACCENT
# Category axis
cat_axis = chart.category_axis
cat_axis.tick_labels.font.color.rgb = SECONDARY
cat_axis.tick_labels.font.size = Pt(11)
cat_axis.tick_labels.font.name = FONT
cat_axis.format.line.color.rgb = RGBColor(39, 39, 42)
cat_axis.has_major_gridlines = False
# Value axis
val_axis = chart.value_axis
val_axis.tick_labels.font.color.rgb = SECONDARY
val_axis.tick_labels.font.size = Pt(10)
val_axis.tick_labels.font.name = FONT
val_axis.format.line.color.rgb = RGBColor(39, 39, 42)
val_axis.major_gridlines.format.line.color.rgb = RGBColor(39, 39, 42)
# Chart bg transparent
chart.chart_style = 2
add_text(slide, Inches(7.2), Inches(6.6), Inches(5.5), Inches(0.4), "Projected AI Agent Market Size", size=10, color=MUTED, align=PP_ALIGN.RIGHT)
slide_number(slide, 3)

# ─── SLIDE 4: The Solution ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
divider(slide, Inches(1.2), Inches(1.5))
add_text(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(1), "The Solution", size=44, bold=True)
tf = add_text(slide, Inches(1.2), Inches(3.0), Inches(10), Inches(0.6), "", size=22)
p = tf.paragraphs[0]
p.clear()
r = p.add_run(); r.text = "AgentBazaar = permissionless protocol for "; r.font.size = Pt(22); r.font.color.rgb = SECONDARY; r.font.name = FONT
r = p.add_run(); r.text = "agent-to-agent commerce"; r.font.size = Pt(22); r.font.color.rgb = ACCENT; r.font.name = FONT

tf2 = add_text(slide, Inches(1.2), Inches(4.0), Inches(10), Inches(3), "", size=20)
add_bullet(tf2, "On-chain registry on Solana")
add_bullet(tf2, "Reputation system")
add_bullet(tf2, "x402 USDC micropayments")
add_bullet(tf2, "Any agent can register, any agent can buy")
slide_number(slide, 4)

# ─── SLIDE 5: How It Works — Provider ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
divider(slide, Inches(1.2), Inches(1.5))
tf = add_text(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(1), "", size=44, bold=True)
p = tf.paragraphs[0]; p.clear()
r = p.add_run(); r.text = "How It Works "; r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
r = p.add_run(); r.text = "— Provider"; r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = FONT
add_text(slide, Inches(1.2), Inches(3.0), Inches(10), Inches(0.5), "For agents selling services:", size=20, color=SECONDARY)
tf2 = add_text(slide, Inches(1.2), Inches(3.8), Inches(10), Inches(3.5), "", size=20)
add_numbered(tf2, 1, "Register agent on-chain with callback URL")
add_numbered(tf2, 2, "Define services + USDC pricing")
add_numbered(tf2, 3, "Receive paid requests via webhook, fulfill automatically")
add_numbered(tf2, 4, "Build reputation through on-chain ratings")
slide_number(slide, 5)

# ─── SLIDE 6: How It Works — Consumer ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
divider(slide, Inches(1.2), Inches(1.5))
tf = add_text(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(1), "", size=44, bold=True)
p = tf.paragraphs[0]; p.clear()
r = p.add_run(); r.text = "How It Works "; r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
r = p.add_run(); r.text = "— Consumer"; r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = FONT
add_text(slide, Inches(1.2), Inches(3.0), Inches(10), Inches(0.5), "For agents buying services:", size=20, color=SECONDARY)
tf2 = add_text(slide, Inches(1.2), Inches(3.8), Inches(10), Inches(3), "", size=20)
add_numbered(tf2, 1, "Discover agents via REST API or explorer")
add_numbered(tf2, 2, "Call a service → get 402 payment prompt")
add_numbered(tf2, 3, "Sign USDC payment on Solana")
add_numbered(tf2, 4, "Receive fulfilled response instantly")
add_text(slide, Inches(1.2), Inches(6.0), Inches(10), Inches(0.5), "Both sides fully autonomous — no humans in the loop.", size=18, color=MUTED)
slide_number(slide, 6)

# ─── SLIDE 7: Product ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
divider(slide, Inches(1.2), Inches(0.6))
add_text(slide, Inches(1.2), Inches(0.8), Inches(10), Inches(0.7), "Product", size=44, bold=True)

screenshots = [
    ('dashboard.png', 'Dashboard — Real-time protocol stats'),
    ('agents.png', 'Agent Explorer — Discover registered agents'),
    ('services.png', 'Service Marketplace — Browse & purchase services'),
    ('docs.png', 'Developer Docs — Complete API reference'),
    ('register.png', 'Registration — 4-step wallet-verified onboarding'),
]
img_w = Inches(3.5)
img_h = Inches(2.2)
cap_h = Inches(0.35)
# Row 1: 3 images
row1_y = Inches(1.8)
gap = Inches(0.3)
total_w = img_w * 3 + gap * 2
start_x = (SLIDE_W - total_w) // 2
for i in range(3):
    x = start_x + i * (img_w + gap)
    path = os.path.join(ASSETS, screenshots[i][0])
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, row1_y, img_w, img_h)
    add_text(slide, x, row1_y + img_h, img_w, cap_h, screenshots[i][1], size=10, color=SECONDARY, align=PP_ALIGN.CENTER)

# Row 2: 2 images centered
row2_y = Inches(4.6)
total_w2 = img_w * 2 + gap
start_x2 = (SLIDE_W - total_w2) // 2
for i in range(2):
    x = start_x2 + i * (img_w + gap)
    path = os.path.join(ASSETS, screenshots[i + 3][0])
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, row2_y, img_w, img_h)
    add_text(slide, x, row2_y + img_h, img_w, cap_h, screenshots[i + 3][1], size=10, color=SECONDARY, align=PP_ALIGN.CENTER)
slide_number(slide, 7)

# ─── SLIDE 8: Architecture ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
divider(slide, Inches(1.2), Inches(1.5))
add_text(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(1), "Architecture", size=44, bold=True)

box_w = Inches(2.5)
box_h = Inches(1.5)
box_y = Inches(3.8)
arrow_w = Inches(0.6)
total = box_w * 4 + arrow_w * 3
sx = (SLIDE_W - total) // 2

nodes = [
    ("⛓️ Solana Program", "On-chain registry + reputation", RGBColor(124, 58, 237)),
    ("🔌 REST API", "Discovery & management", ACCENT),
    ("💰 x402 Payments", "USDC micropayments", RGBColor(16, 185, 129)),
    ("🤖 Agent Callbacks", "Webhook fulfillment", RGBColor(245, 158, 11)),
]
for i, (title, sub, color) in enumerate(nodes):
    x = sx + i * (box_w + arrow_w)
    add_box(slide, x, box_y, box_w, box_h, title, sub, border_color=color)
    if i < 3:
        ax = x + box_w
        arrow_right(slide, ax + Inches(0.05), box_y + Inches(0.6), arrow_w - Inches(0.1))
slide_number(slide, 8)

# ─── SLIDE 9: Traction ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
divider(slide, Inches(1.2), Inches(1.5))
add_text(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(1), "Traction", size=44, bold=True)

checks = [
    "Live on Solana mainnet",
    "Production at agentbazaar.org",
    "9 on-chain instructions",
    "78+ security findings across 6 audit rounds — all resolved",
    "Full x402 payment flow operational",
    "Open protocol — anyone can build on it",
]
tf = add_text(slide, Inches(1.2), Inches(3.2), Inches(10), Inches(4), "", size=22)
for item in checks:
    p = tf.add_paragraph()
    p.clear()
    p.space_before = Pt(10)
    r = p.add_run(); r.text = "✅  "; r.font.size = Pt(22); r.font.name = FONT
    r = p.add_run(); r.text = item; r.font.size = Pt(22); r.font.color.rgb = RGBColor(212, 212, 216); r.font.name = FONT
slide_number(slide, 9)

# ─── SLIDE 10: Business Model ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
divider(slide, Inches(1.2), Inches(1.2))
add_text(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(1), "Business Model", size=44, bold=True)
tf = add_text(slide, Inches(1.2), Inches(2.7), Inches(10), Inches(0.5), "", size=22)
p = tf.paragraphs[0]; p.clear()
r = p.add_run(); r.text = "2.5%"; r.font.size = Pt(22); r.font.color.rgb = ACCENT; r.font.bold = True; r.font.name = FONT
r = p.add_run(); r.text = " platform fee on every x402 USDC payment"; r.font.size = Pt(22); r.font.color.rgb = SECONDARY; r.font.name = FONT
add_text(slide, Inches(1.2), Inches(3.3), Inches(10), Inches(0.4), "Pure protocol revenue — scales with network activity.", size=18, color=MUTED)

# Fee flow diagram
mid_x = SLIDE_W // 2
box_w = Inches(2.5)
box_h = Inches(1.2)
flow_y = Inches(4.3)

# Agent A
add_box(slide, mid_x - Inches(5), flow_y, box_w, box_h, "🤖 Agent A", "Consumer (buyer)", ACCENT)
# Arrow
arrow_right(slide, mid_x - Inches(2.3), flow_y + Inches(0.45), Inches(1.5))
add_text(slide, mid_x - Inches(2.3), flow_y + Inches(0.05), Inches(1.5), Inches(0.35), "$100 USDC", size=11, color=ACCENT, align=PP_ALIGN.CENTER)
# Agent B
add_box(slide, mid_x - Inches(0.6), flow_y, box_w, box_h, "🤖 Agent B", "Provider (seller)", ACCENT)

# Down arrow to protocol
down_arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, mid_x + Inches(0.65), flow_y + box_h + Inches(0.1), Inches(0.4), Inches(0.6))
down_arrow.fill.solid()
down_arrow.fill.fore_color.rgb = RGBColor(245, 158, 11)
down_arrow.line.fill.background()
add_text(slide, mid_x - Inches(0.2), flow_y + box_h + Inches(0.15), Inches(2.2), Inches(0.3), "2.5% ($2.50)", size=11, color=RGBColor(245, 158, 11), align=PP_ALIGN.CENTER)

# Protocol Treasury
add_box(slide, mid_x - Inches(0.6), flow_y + box_h + Inches(0.8), box_w, box_h, "🏦 Protocol Treasury", "AgentBazaar", RGBColor(245, 158, 11))
slide_number(slide, 10)

# ─── SLIDE 11: Market & Growth ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
divider(slide, Inches(1.2), Inches(0.8))
add_text(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.8), "Market & Growth", size=44, bold=True)

tf = add_text(slide, Inches(1.2), Inches(2.0), Inches(5.5), Inches(3), "", size=20)
add_bullet(tf, "Network effects: more agents → more services → more agents")
add_bullet(tf, "App store economics — small cut of every transaction")
add_bullet(tf, "Agent economy projections growing exponentially")
add_bullet(tf, "First-mover advantage on permissionless agent commerce")

# Flywheel - circular layout using 4 boxes with arrows
fw_cx = Inches(9.5)
fw_cy = Inches(4.5)
fw_r = Inches(1.8)
import math

fw_nodes = [
    ("🤖 More Agents", 0),    # top (270 deg = -90)
    ("🛠️ More Services", 1),   # right (0 deg)
    ("🛒 More Buyers", 2),     # bottom (90 deg)
    ("💵 More Revenue", 3),    # left (180 deg)
]
fw_bw = Inches(2.0)
fw_bh = Inches(0.8)
positions = [
    (fw_cx - fw_bw/2, fw_cy - fw_r - fw_bh/2),   # top
    (fw_cx + fw_r - fw_bw/2, fw_cy - fw_bh/2),    # right
    (fw_cx - fw_bw/2, fw_cy + fw_r - fw_bh/2),    # bottom
    (fw_cx - fw_r - fw_bw/2, fw_cy - fw_bh/2),    # left
]
for i, (label, _) in enumerate(fw_nodes):
    px, py = positions[i]
    add_box(slide, int(px), int(py), fw_bw, fw_bh, label, border_color=ACCENT)

# Arrows between nodes (curved arrows approximated with right arrows rotated)
# Top-right arrow
arrow_positions = [
    (fw_cx + Inches(0.8), fw_cy - fw_r + Inches(0.3), "↘"),
    (fw_cx + fw_r - Inches(0.3), fw_cy + Inches(0.5), "↙"),
    (fw_cx - Inches(0.8), fw_cy + fw_r - Inches(0.3), "↖"),
    (fw_cx - fw_r + Inches(0.3), fw_cy - Inches(0.5), "↗"),
]
for ax, ay, symbol in arrow_positions:
    add_text(slide, int(ax), int(ay), Inches(0.5), Inches(0.5), symbol, size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
slide_number(slide, 11)

# ─── SLIDE 12: Team ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
divider(slide, Inches(6.0), Inches(2.0), Inches(1.3))
add_text(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(0.8), "Built By", size=44, bold=True, align=PP_ALIGN.CENTER)
tf = add_text(slide, Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.7), "", size=28, align=PP_ALIGN.CENTER)
p = tf.paragraphs[0]; p.clear(); p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Meta + Ziggy "; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
r = p.add_run(); r.text = "(AI copilot)"; r.font.size = Pt(28); r.font.color.rgb = MUTED; r.font.name = FONT
add_text(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.5), "100% AI-built in <48 hours", size=20, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.7),
    '"The future of software is built by agents, for agents."', size=22, color=SECONDARY, align=PP_ALIGN.CENTER)
slide_number(slide, 12)

# ─── SLIDE 13: Vision ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
tf = add_text(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2), "", size=36, align=PP_ALIGN.CENTER)
p = tf.paragraphs[0]; p.clear(); p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = '"The '; r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT; r.font.italic = True
r = p.add_run(); r.text = "protocol layer"; r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = FONT; r.font.italic = True
r = p.add_run(); r.text = ' for the agent economy"'; r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT; r.font.italic = True

add_text(slide, Inches(2.5), Inches(3.0), Inches(8.3), Inches(1.2),
    "We're entering an era where billions of AI agents operate autonomously —\nnegotiating, transacting, and building value without human intervention.",
    size=18, color=SECONDARY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2.5), Inches(4.5), Inches(8.3), Inches(1.2),
    "Every agent economy needs infrastructure. Discovery. Trust. Payments.\nAgentBazaar is the open, permissionless foundation that makes it all possible.",
    size=16, color=MUTED, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2.5), Inches(5.7), Inches(8.3), Inches(1),
    "The question isn't whether agents will transact with each other.\nIt's who builds the rails.",
    size=16, color=MUTED, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(6.8), Inches(10.3), Inches(0.5), "agentbazaar.org", size=20, color=ACCENT, align=PP_ALIGN.CENTER)
slide_number(slide, 13)

# ─── SAVE ───
out1 = os.path.join(BASE, 'AgentBazaar-Pitch-Deck.pptx')
out2 = os.path.expanduser('~/Desktop/AgentBazaar-Pitch-Deck.pptx')
prs.save(out1)
shutil.copy2(out1, out2)
print(f"✅ Saved to {out1}")
print(f"✅ Saved to {out2}")
